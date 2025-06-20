import csv
import logging
import os
from typing import Dict, Callable, Optional
from pathlib import Path
import xml.etree.ElementTree as ET
from bs4 import BeautifulSoup
from striprtf.striprtf import rtf_to_text
from dbfread import DBF

# Document processing libraries (conditional imports)
try:
    import docx
    from pptx import Presentation
    from PyPDF2 import PdfReader
    import openpyxl
    import xlrd
    from odf import text, teletype
    from odf.opendocument import load
    from odf.text import P
    from goblintools.ocr_parser import OCRProcessor
except ImportError as e:
    logging.warning(f"Optional dependency not found: {e}. Some file formats may not be supported.")

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

class TextExtractor:
    """Main class for handling text extraction from various file formats."""
    
    def __init__(self, ocr_handler = False, use_aws=False, aws_access_key=None, aws_secret_key=None, aws_region='us-east-1'):
        """
        Initialize the text extractor.
        
        Args:
            ocr_handler: Optional function to handle OCR for image-based PDFs
        """
        if ocr_handler:
            self.ocr_handler = OCRProcessor(use_aws, aws_access_key, aws_secret_key, aws_region)
        else:
            self.ocr_handler = None

        self._parsers = self._initialize_parsers()

    def _initialize_parsers(self) -> Dict[str, Callable]:
        """Initialize all available text extraction parsers."""
        return {
            '.pdf': self._extract_pdf,
            '.doc': self._extract_doc,
            '.docx': self._extract_docx,
            '.txt': self._extract_txt,
            '.pptx': self._extract_pptx,
            '.html': self._extract_html,
            '.odt': self._extract_odt,
            '.rtf': self._extract_rtf,
            '.csv': self._extract_csv,
            '.xml': self._extract_xml,
            '.xlsx': self._extract_xlsx,
            '.xlsm': self._extract_xlsx,
            '.xls': self._extract_xls,
            '.ods': self._extract_ods,
            '.dbf': self._extract_dbf,
        }

    def add_parser(self, extension: str, parser_func: Callable) -> None:
        """Add or override a parser for a specific file extension."""
        self._parsers[extension.lower()] = parser_func

    def extract_from_file(self, file_path: str) -> str:
        """
        Extract text from a single file.
        
        Args:
            file_path: Path to the file to extract text from
            
        Returns:
            Extracted text as a string
        """
        if not os.path.exists(file_path):
            logger.warning(f"File not found: {file_path}")
            return ""

        file_extension = Path(file_path).suffix.lower()
        parser = self._parsers.get(file_extension)
        
        if not parser:
            logger.warning(f"No parser available for file extension: {file_extension}")
            return ""

        try:
            return parser(file_path)
        except Exception as e:
            logger.error(f"Error extracting text from {file_path}: {e}")
            return ""

    def extract_from_folder(self, folder_path: str) -> str:
        """
        Extract text from all supported files in a folder (recursively).
        
        Args:
            folder_path: Path to the folder to process
            
        Returns:
            Combined extracted text from all files
        """
        if not os.path.exists(folder_path):
            logger.warning(f"Folder not found: {folder_path}")
            return ""

        extracted_text = []
        for root, _, files in os.walk(folder_path):
            for file in files:
                file_path = os.path.join(root, file)
                text = self.extract_from_file(file_path)
                if text:
                    extracted_text.append(text)

        return ' '.join(extracted_text)

    # Individual parser methods
    def _extract_pdf(self, file_path: str) -> str:
        """Extract text from PDF files, with fallback to OCR if needed."""
        has_images = False
        extracted_text = ''
        
        with open(file_path, 'rb') as f:
            reader = PdfReader(f)
            for page in reader.pages:
                extracted_text += page.extract_text() or ''
                if not has_images and '/XObject' in page.get('/Resources', {}):
                    xObject = page['/Resources']['/XObject'].get_object()
                    has_images = any(
                        xObject[obj]['/Subtype'] == '/Image'
                        for obj in xObject
                    )

        # Fallback to OCR if text extraction failed and images were found
        if (not extracted_text.strip()) and has_images and self.ocr_handler:
            return self.ocr_handler.extract_text_from_pdf(file_path)
        elif not self.ocr_handler:
            logger.warning(f"The file {file_path}, requires OCR but it was not enabled.")
        else:
            logger.warning(f"Unable to process {file_path}.")
            
        return extracted_text

    def _extract_docx(self, file_path: str) -> str:
        """Extract text from DOCX files."""
        try:
            doc = docx.Document(file_path)
            return ' '.join(para.text for para in doc.paragraphs if para.text)
        except Exception as e:
            logger.error(f"Error processing DOCX file {file_path}: {e}")
            return ""

    def _extract_doc(self, file_path: str) -> str:
        """Extract text from legacy DOC files."""
        try:
            return textract.process(file_path).decode('utf-8')
        except Exception as e:
            logger.error(f"Error processing DOC file {file_path}: {e}")
            return ""

    def _extract_txt(self, file_path: str) -> str:
        """Extract text from plain text files."""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                return file.read()
        except UnicodeDecodeError:
            try:
                with open(file_path, 'r', encoding='latin-1') as file:
                    return file.read()
            except Exception as e:
                logger.error(f"Error processing TXT file {file_path}: {e}")
                return ""

    def _extract_pptx(self, file_path: str) -> str:
        """Extract text from PowerPoint files."""
        try:
            prs = Presentation(file_path)
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        texts.append(shape.text)
            return ' '.join(texts)
        except Exception as e:
            logger.error(f"Error processing PPTX file {file_path}: {e}")
            return ""

    def _extract_html(self, file_path: str) -> str:
        """Extract text from HTML files."""
        encodings = ['utf-8', 'latin-1']
        for encoding in encodings:
            try:
                with open(file_path, 'r', encoding=encoding) as file:
                    soup = BeautifulSoup(file.read(), 'html.parser')
                    return soup.get_text(separator=' ', strip=True)
            except UnicodeDecodeError:
                continue
            except Exception as e:
                logger.error(f"Error processing HTML file {file_path}: {e}")
                break
        return ""

    def _extract_odt(self, file_path: str) -> str:
        """Extract text from OpenDocument Text files."""
        try:
            doc = load(file_path)
            return ' '.join(
                teletype.extractText(element)
                for element in doc.getElementsByType(text.P)
            )
        except Exception as e:
            logger.error(f"Error processing ODT file {file_path}: {e}")
            return ""

    def _extract_rtf(self, file_path: str) -> str:
        """Extract text from RTF files."""
        try:
            with open(file_path, 'r') as file:
                return rtf_to_text(file.read(), errors='ignore')
        except Exception as e:
            logger.error(f"Error processing RTF file {file_path}: {e}")
            return ""

    def _extract_csv(self, file_path: str) -> str:
        """Extract text from CSV files."""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                return ' '.join(
                    ' '.join(row) 
                    for row in csv.reader(file)
                    if any(field.strip() for field in row)
                )
        except Exception as e:
            logger.error(f"Error processing CSV file {file_path}: {e}")
            return ""

    def _extract_xml(self, file_path: str) -> str:
        """Extract text from XML files."""
        try:
            tree = ET.parse(file_path)
            return ' '.join(
                elem.text.strip() 
                for elem in tree.iter() 
                if elem.text and elem.text.strip()
            )
        except Exception as e:
            logger.error(f"Error processing XML file {file_path}: {e}")
            return ""

    def _extract_xlsx(self, file_path: str) -> str:
        """Extract text from modern Excel files."""
        try:
            wb = openpyxl.load_workbook(file_path)
            texts = []
            for sheet in wb:
                texts.extend(
                    str(cell) 
                    for row in sheet.iter_rows(values_only=True) 
                    for cell in row 
                    if cell is not None
                )
            return ' '.join(texts)
        except Exception as e:
            logger.error(f"Error processing XLSX file {file_path}: {e}")
            return ""

    def _extract_xls(self, file_path: str) -> str:
        """Extract text from legacy Excel files."""
        try:
            sheet = xlrd.open_workbook(file_path).sheet_by_index(0)
            return ' '.join(
                str(col.value) 
                for rw in range(sheet.nrows) 
                for col in sheet.row(rw) 
                if col.value
            )
        except Exception as e:
            logger.error(f"Error processing XLS file {file_path}: {e}")
            return ""

    def _extract_ods(self, file_path: str) -> str:
        """Extract text from OpenDocument Spreadsheets."""
        try:
            doc = load(file_path)
            return '\n'.join(
                "".join(
                    child.data 
                    for child in p.childNodes 
                    if child.nodeType == child.TEXT_NODE
                )
                for p in doc.getElementsByType(P)
            )
        except Exception as e:
            logger.error(f"Error processing ODS file {file_path}: {e}")
            return ""

    def _extract_dbf(self, file_path: str) -> str:
        """Extract text from DBF database files."""
        try:
            dbf = DBF(file_path, load=True)
            return ' '.join(
                f"{key}: {value}" 
                for record in dbf 
                for key, value in record.items()
            )
        except Exception as e:
            logger.error(f"Error processing DBF file {file_path}: {e}")
            return ""


# Legacy functions for backward compatibility
def parsers() -> Dict[str, Callable]:
    extractor = TextExtractor()
    return extractor._parsers

def parse_to_string(file_path: str) -> str:
    return TextExtractor().extract_from_file(file_path)

def files_to_string(folder_path: str) -> str:
    return TextExtractor().extract_from_folder(folder_path)
